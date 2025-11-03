"""
Multi-format file processor for PDFs, images, documents, etc.
"""
from pathlib import Path
from typing import List, Union
from PIL import Image
import io
import logging

logger = logging.getLogger(__name__)


def process_file(file_path: str, embedding_client) -> List[float]:
    """
    Process any supported file type and return embeddings.
    
    Args:
        file_path: Path to the file
        embedding_client: Client for embedding generation
    
    Returns:
        List of embeddings (vector)
    """
    file_path = Path(file_path)
    extension = file_path.suffix.lower()
    
    if extension == '.pdf':
        return process_pdf(file_path, embedding_client)
    elif extension in ['.png', '.jpg', '.jpeg']:
        return process_image(file_path, embedding_client)
    elif extension in ['.docx']:
        return process_docx(file_path, embedding_client)
    elif extension in ['.pptx']:
        return process_pptx(file_path, embedding_client)
    elif extension in ['.xlsx', '.csv']:
        return process_excel(file_path, embedding_client)
    elif extension in ['.txt', '.md']:
        return process_text(file_path, embedding_client)
    else:
        raise ValueError(f"Unsupported file type: {extension}")


def process_pdf(file_path: Path, embedding_client) -> List[float]:
    """Process PDF using multimodal pipeline."""
    from semantic.parsing_engine import parse_local_pdf
    from semantic.embedding_engine import embed_pdf
    
    images = parse_local_pdf(str(file_path))
    return embed_pdf(images, embedding_client)


def process_image(file_path: Path, embedding_client) -> List[float]:
    """Process image with OCR then text embedding."""
    try:
        from PIL import Image
        
        # Open image
        image = Image.open(file_path)
        
        # Try OCR if available
        try:
            import pytesseract
            text = pytesseract.image_to_string(image)
            
            if text.strip():
                # OCR successful, use text embedding
                return embed_text(text, embedding_client)
        except ImportError:
            logger.info("Tesseract not installed, using multimodal embedding for image")
        except Exception as e:
            logger.warning(f"OCR failed for {file_path}, using multimodal: {e}")
        
        # Fallback: use image directly with multimodal embedding
        if hasattr(embedding_client, 'multimodal_embed'):
            return embedding_client.multimodal_embed(
                inputs=[[image]],
                model="voyage-multimodal-3",
                input_type="document"
            ).embeddings[0]
        else:
            raise ValueError("Embedding client does not support multimodal embedding")
        
    except Exception as e:
        logger.error(f"Failed to process image {file_path}: {e}")
        raise


def process_docx(file_path: Path, embedding_client) -> List[float]:
    """Process Word document."""
    try:
        from docx import Document
        
        doc = Document(file_path)
        
        # Extract all text from paragraphs
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_parts.append(cell.text)
        
        full_text = '\n'.join(text_parts)
        
        if not full_text.strip():
            raise ValueError("No text content found in document")
        
        return embed_text(full_text, embedding_client)
        
    except Exception as e:
        logger.error(f"Failed to process DOCX {file_path}: {e}")
        raise


def process_pptx(file_path: Path, embedding_client) -> List[float]:
    """Process PowerPoint presentation."""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text_parts = []
        
        # Extract text from all slides
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_parts.append(shape.text)
        
        full_text = '\n'.join(text_parts)
        
        if not full_text.strip():
            raise ValueError("No text content found in presentation")
        
        return embed_text(full_text, embedding_client)
        
    except Exception as e:
        logger.error(f"Failed to process PPTX {file_path}: {e}")
        raise


def process_excel(file_path: Path, embedding_client) -> List[float]:
    """Process Excel spreadsheet."""
    try:
        from openpyxl import load_workbook
        
        wb = load_workbook(file_path, read_only=True, data_only=True)
        text_parts = []
        
        # Extract text from all sheets
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = ' '.join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    text_parts.append(row_text)
        
        full_text = '\n'.join(text_parts)
        
        if not full_text.strip():
            raise ValueError("No content found in spreadsheet")
        
        return embed_text(full_text, embedding_client)
        
    except Exception as e:
        logger.error(f"Failed to process Excel {file_path}: {e}")
        raise


def process_text(file_path: Path, embedding_client) -> List[float]:
    """Process plain text or markdown file with encoding fallback."""
    try:
        # Try UTF-8 first
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
        except UnicodeDecodeError:
            # Fallback to latin-1 for non-UTF8 files
            with open(file_path, 'r', encoding='latin-1', errors='ignore') as f:
                text = f.read()
        
        if not text.strip():
            raise ValueError("Empty file")
        
        return embed_text(text, embedding_client)
        
    except Exception as e:
        logger.error(f"Failed to process text file {file_path}: {e}")
        raise


def embed_text(text: str, embedding_client) -> List[float]:
    """
    Embed text using the provided client.
    
    Args:
        text: Text to embed (truncate if too long)
        embedding_client: Client for embedding generation
    
    Returns:
        Embedding vector
    """
    # Truncate if needed
    MAX_CHARS = 32000
    if len(text) > MAX_CHARS:
        text = text[:MAX_CHARS]
        logger.warning(f"Text truncated to {MAX_CHARS} characters")
    
    try:
        if hasattr(embedding_client, 'embed'):
            response = embedding_client.embed(
                texts=[text],
                model="voyage-3",
                input_type="document"
            )
            return response.embeddings[0]
        else:
            raise ValueError("Embedding client does not support text embedding")
    except Exception as e:
        logger.error(f"Failed to embed text: {e}")
        raise


def get_supported_extensions() -> List[str]:
    """Get list of supported file extensions."""
    return [
        '.pdf',
        '.png', '.jpg', '.jpeg', '.gif', '.webp',
        '.docx', '.doc',
        '.pptx', '.ppt',
        '.xlsx', '.xls', '.csv',
        '.txt', '.md', '.rtf'
    ]


def is_supported_file(file_path: Union[str, Path]) -> bool:
    """Check if file type is supported and not a system file."""
    path = Path(file_path)
    
    # Skip system/thumbnail files
    skip_extensions = ['.thm', '.tmp', '.cache', '.ds_store']
    if path.suffix.lower() in skip_extensions:
        return False
    
    # Skip hidden files and system folders
    if path.name.startswith('.') or '/.git/' in str(path) or '/node_modules/' in str(path):
        return False
    
    return path.suffix.lower() in get_supported_extensions()

